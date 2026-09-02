"""문서 파싱: PDF / PPTX / TXT 를 페이지별 텍스트로 바꾼다 (docs/02-document-analysis.md).

오류는 사용자에게 그대로 보여줄 수 있는 한국어 메시지를 담는다 (docs/10-quality-safety.md).
"""

from __future__ import annotations

import io
from dataclasses import dataclass
from pathlib import Path

SUPPORTED_SUFFIXES = {".pdf", ".pptx", ".txt", ".md", ".text"}

# 형식별 안내가 필요한 미지원 확장자 — 무엇으로 바꿔 오면 되는지까지 알려준다.
_CONVERTIBLE_SUFFIXES = {
    ".ppt": "구버전 PowerPoint(.ppt) 파일은 지원하지 않습니다. PowerPoint에서 .pptx 로 저장한 뒤 다시 올려 주세요.",
    ".doc": "Word 문서(.doc)는 지원하지 않습니다. PDF로 저장한 뒤 다시 올려 주세요.",
    ".docx": "Word 문서(.docx)는 지원하지 않습니다. PDF로 저장한 뒤 다시 올려 주세요.",
    ".hwp": "한글 문서(.hwp)는 지원하지 않습니다. PDF로 저장한 뒤 다시 올려 주세요.",
    ".key": "Keynote 파일은 지원하지 않습니다. PDF 또는 .pptx 로 내보낸 뒤 다시 올려 주세요.",
}


class DocumentError(ValueError):
    """사용자에게 그대로 노출해도 되는 문서 처리 오류."""


@dataclass
class PageText:
    page: int
    text: str


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp949", "euc-kr"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentError("텍스트 인코딩을 인식하지 못했습니다. UTF-8로 저장한 뒤 다시 시도해 주세요.")


def _parse_pdf(data: bytes) -> list[PageText]:
    try:
        # 구버전 별칭 `fitz` 는 deprecation 경고를 낸다. 새 이름을 먼저 시도한다.
        import pymupdf
    except ImportError:
        try:
            import fitz as pymupdf  # type: ignore[no-redef]
        except ImportError as exc:  # pragma: no cover - 설치 환경 문제
            raise DocumentError(
                "PDF 처리 모듈(PyMuPDF)이 설치되어 있지 않습니다. TXT 파일을 사용하거나 "
                "requirements.txt 를 설치해 주세요."
            ) from exc

    try:
        document = pymupdf.open(stream=data, filetype="pdf")
    except Exception as exc:  # noqa: BLE001 - 손상 파일 등
        raise DocumentError(f"PDF를 열 수 없습니다: {exc}") from exc

    pages: list[PageText] = []
    with document:
        for index, page in enumerate(document, start=1):
            pages.append(PageText(page=index, text=page.get_text("text") or ""))
    return pages


def _shape_text(shape) -> str:  # noqa: ANN001 - python-pptx 타입은 런타임에만 있다
    """도형 하나에서 근거로 쓸 텍스트를 뽑는다. 표는 행 단위로 편다."""
    if getattr(shape, "has_table", False):
        rows: list[str] = []
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            # 병합 셀은 같은 텍스트가 반복되므로 연속 중복만 접는다
            merged: list[str] = []
            for cell in cells:
                if not merged or merged[-1] != cell:
                    merged.append(cell)
            line = " | ".join(merged).strip(" |")
            if line:
                rows.append(line)
        return "\n".join(rows)

    if getattr(shape, "has_text_frame", False):
        lines = [p.text.strip() for p in shape.text_frame.paragraphs]
        return "\n".join(line for line in lines if line)

    return ""


def _iter_shape_texts(shapes) -> list[str]:  # noqa: ANN001
    """도형 목록을 읽기 순서(위→아래, 왼→오른쪽)로 훑어 텍스트 블록을 모은다.

    z-order 는 화면 배치와 무관해서 그대로 읽으면 제목이 본문 뒤로 가기도 한다.
    그룹 도형은 안쪽까지 내려간다.
    """
    ordered = sorted(shapes, key=lambda s: (getattr(s, "top", None) or 0, getattr(s, "left", None) or 0))

    blocks: list[str] = []
    for shape in ordered:
        if hasattr(shape, "shapes"):  # 그룹 도형
            blocks.extend(_iter_shape_texts(shape.shapes))
            continue
        text = _shape_text(shape)
        if text.strip():
            blocks.append(text.strip())
    return blocks


def _parse_pptx(data: bytes) -> list[PageText]:
    """슬라이드 1장을 페이지 1개로 본다. 발표자 노트도 원문의 일부로 함께 읽는다."""
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - 설치 환경 문제
        raise DocumentError(
            "PPTX 처리 모듈(python-pptx)이 설치되어 있지 않습니다. PDF나 TXT 파일을 사용하거나 "
            "requirements.txt 를 설치해 주세요."
        ) from exc

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 - 손상 파일, .ppt 를 확장자만 바꾼 파일 등
        raise DocumentError(
            f"PPTX를 열 수 없습니다: {exc} "
            "구버전 .ppt 파일이라면 PowerPoint에서 .pptx 로 저장한 뒤 다시 시도해 주세요."
        ) from exc

    pages: list[PageText] = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []

        # 제목은 배치와 무관하게 맨 앞에 둔다 (슬라이드의 주제를 chunk 앞머리에 남기기 위해)
        title = None
        try:
            title = slide.shapes.title
        except (AttributeError, ValueError):  # pragma: no cover - 레이아웃에 제목이 없는 경우
            title = None
        title_text = _shape_text(title).strip() if title is not None else ""
        if title_text:
            blocks.append(title_text)

        for text in _iter_shape_texts(slide.shapes):
            if text != title_text:
                blocks.append(text)

        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                blocks.append(f"[발표자 노트] {notes}")

        pages.append(PageText(page=index, text="\n\n".join(blocks)))

    if not pages:
        raise DocumentError("슬라이드가 없는 PPTX 파일입니다. 내용이 있는 파일을 올려 주세요.")

    return pages


def _split_plain_text_pages(text: str, chars_per_page: int = 1200) -> list[PageText]:
    """TXT 에는 페이지 개념이 없다. 근거 표시를 위해 일정 길이로 가상 페이지를 만든다."""
    blocks = [b for b in text.split("\n\n") if b.strip()]
    pages: list[PageText] = []
    current: list[str] = []
    length = 0

    for block in blocks:
        if current and length + len(block) > chars_per_page:
            pages.append(PageText(page=len(pages) + 1, text="\n\n".join(current)))
            current, length = [], 0
        current.append(block)
        length += len(block)

    if current:
        pages.append(PageText(page=len(pages) + 1, text="\n\n".join(current)))
    return pages or [PageText(page=1, text=text)]


def parse_document(filename: str, data: bytes) -> list[PageText]:
    """업로드된 파일을 페이지별 텍스트로 변환한다."""
    suffix = Path(filename or "").suffix.lower()

    if suffix not in SUPPORTED_SUFFIXES:
        hint = _CONVERTIBLE_SUFFIXES.get(suffix)
        raise DocumentError(
            hint
            or (
                "PDF, PPTX, TXT 파일만 업로드할 수 있습니다. "
                f"(받은 형식: {suffix or '알 수 없음'})"
            )
        )

    if not data:
        raise DocumentError("업로드된 파일이 비어 있습니다.")

    if suffix == ".pdf":
        pages = _parse_pdf(data)
    elif suffix == ".pptx":
        pages = _parse_pptx(data)
    else:
        pages = _split_plain_text_pages(_decode_text(data))

    if not any(page.text.strip() for page in pages):
        if suffix == ".pptx":
            raise DocumentError(
                "슬라이드에서 텍스트를 찾지 못했습니다. 이미지로만 만든 발표자료일 수 있습니다. "
                "텍스트 상자가 있는 파일로 다시 시도해 주세요."
            )
        raise DocumentError(
            "문서에서 텍스트를 찾지 못했습니다. 스캔 이미지 PDF일 수 있습니다. "
            "텍스트가 포함된 파일로 다시 시도해 주세요."
        )

    return pages
