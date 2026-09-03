"""SlideDeck -> PPTX 파일 (docs/04-slide-planner.md 의 선택 항목).

이 모듈은 **새로운 문장을 만들지 않는다.** 이미 검증을 마친 `GenerateResponse` 를 그대로
슬라이드 위에 배치할 뿐이다. 여기서 문장을 요약하거나 자르면 원문 근거와의 대응이 깨진다.

입력이 PPTX 였으면 원본 파일 위에 결과를 얹는다. 원본의 이미지·표·배경·글꼴은 그대로 두고
텍스트만 청중용으로 바꾸므로 사용자가 만들어 둔 형식이 살아남는다. 원본이 없거나(PDF·TXT)
얹기에 실패하면 지금까지처럼 빈 프레젠테이션에 새 덱을 그린다.

python-pptx 가 없어도 앱은 뜬다. JSON / Markdown 다운로드만으로 데모가 성립해야 하므로
(docs/04-slide-planner.md) import 실패는 이 모듈 안에 가둬 두고 API 가 한국어로 안내한다.
"""

from __future__ import annotations

import copy
import io
import logging
from collections import Counter

from app.models.contracts import GenerateResponse, Slide
from app.services import labels

try:  # pragma: no cover - 설치 환경에서는 항상 성공한다
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover - 의존성이 빠진 환경
    PPTX_AVAILABLE = False

_log = logging.getLogger(__name__)
_EMU_PER_INCH = 914400


class PptxUnavailableError(RuntimeError):
    """python-pptx 가 설치되어 있지 않다."""


# 화면(frontend/app/globals.css)의 라이트 테마 색을 그대로 쓴다 --------------
_FOREGROUND = "1B1B1A"
_MUTED = "6B6B66"
_BORDER = "E0E0DC"
_ACCENT = "D1600A"
_ACCENT_SOFT = "FDF1E6"
_WARN = "8A5A00"
_WARN_SOFT = "FDF3E0"
_SURFACE_MUTED = "F1F1EF"

# 한글이 깨지지 않는 폰트. 없는 환경에서는 PowerPoint 가 대체 폰트를 쓴다.
_FONT = "맑은 고딕"

# 16:9 — 원본 없이 새로 그릴 때의 기본값. 원본 위에 얹을 때는 원본 크기를 따른다.
_SLIDE_W = 13.333
_SLIDE_H = 7.5
_MARGIN = 0.7

_DISCLAIMER = "이 자료는 업로드한 원문을 근거로 생성되었습니다. 발표 전 담당자 검토가 필요합니다."


# --------------------------------------------------------------------------
# 저수준 도우미
# --------------------------------------------------------------------------


def _rgb(hex_value: str) -> "RGBColor":
    return RGBColor.from_string(hex_value)


def _style_run(run, size: float, *, bold: bool = False, color: str = _FOREGROUND) -> None:
    """글꼴을 지정한다.

    python-pptx 의 `font.name` 은 라틴 문자 글꼴(a:latin)만 설정한다. 한글은 동아시아
    글꼴(a:ea)을 따르므로 둘 다 지정해야 PowerPoint 에서 같은 글꼴로 보인다.
    """
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = _FONT

    rPr = run._r.get_or_add_rPr()  # noqa: SLF001 - 동아시아 글꼴은 공개 API 가 없다
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rPr.insert_element_before(
            ea, "a:cs", "a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst"
        )
    ea.set("typeface", _FONT)


def _textbox(slide, left: float, top: float, width: float, height: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return frame


def _line(frame, text: str, size: float, *, first: bool, bold: bool = False,
          color: str = _FOREGROUND, space_after: float = 6, align=None):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.space_after = Pt(space_after)
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    _style_run(run, size, bold=bold, color=color)
    return paragraph


def _rect(slide, left: float, top: float, width: float, height: float, color: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False  # 기본 그림자는 업무 자료 톤에 맞지 않는다
    return shape


def _geom(presentation) -> tuple[float, float]:
    """슬라이드 폭·높이(인치). 원본 위에 얹을 때는 4:3 일 수도 있어 고정할 수 없다."""
    return presentation.slide_width.inches, presentation.slide_height.inches


def _blank(presentation):
    """빈 슬라이드. 기본 템플릿은 6번 레이아웃이지만 업로드한 원본에는 없을 수 있다."""
    layout = min(presentation.slide_layouts, key=lambda item: len(item.placeholders))
    slide = presentation.slides.add_slide(layout)
    for placeholder in list(slide.placeholders):
        placeholder._element.getparent().remove(placeholder._element)  # noqa: SLF001
    return slide


def _body_size(texts: list[str], base: float = 18) -> float:
    """본문 bullet 글자 크기.

    문장을 자르지 않는 것이 원칙이므로(docs/04-slide-planner.md) 넘칠 것 같으면 글자를 줄인다.
    """
    total = sum(len(text) for text in texts)
    if total > 360 or len(texts) > 6:
        return base - 5
    if total > 220:
        return base - 3
    return base


def _text_height(lines: list[str], size: float, width: float, *, space_after: float) -> float:
    """줄바꿈까지 셈한 대략의 글 높이(인치).

    한글은 글자 하나가 거의 정사각형이라 글자 폭을 글자 크기와 같다고 본다. 라틴 문자는
    이보다 좁아 실제보다 넉넉히 잡히는데, 모자라게 잡아 겹치는 것보다 낫다.
    """
    per_line = max(int(width * 72 / size), 1)
    rows = sum(max(-(-len(line) // per_line), 1) for line in lines)
    return (rows * size * 1.25 + len(lines) * space_after) / 72


def _fit_body_size(
    lines: list[str], width: float, height: float, *, space_after: float, floor: float = 9
) -> float | None:
    """주어진 자리에 들어가는 가장 큰 글자 크기. `floor` 로도 넘치면 None.

    `_body_size` 는 글자 수만 보므로 자리가 좁으면 그대로 넘쳐 아래 footer 를 덮는다.
    문장을 자르지 않는 것이 원칙이라(docs/04-slide-planner.md) 줄일 수 있는 건 글자 크기뿐이고,
    그것으로도 안 되면 호출부가 그 슬라이드를 포기하거나 바닥 크기로 놓는다.
    """
    size = _body_size(lines)
    while _text_height(lines, size, width, space_after=space_after) > height:
        size -= 1
        if size < floor:
            return None
    return size


def _shrink(total_chars: int, *, soft: int, hard: int) -> float:
    """부록처럼 글자 크기가 여러 단계인 영역에서 한꺼번에 줄일 폭."""
    if total_chars > hard:
        return 3.5
    if total_chars > soft:
        return 2.0
    return 0.0


# --------------------------------------------------------------------------
# 슬라이드
# --------------------------------------------------------------------------


def _slide_header(slide, title: str, position: str, w: float) -> None:
    content_w = w - _MARGIN * 2
    frame = _textbox(slide, _MARGIN, 0.5, content_w - 1.7, 0.9)
    _line(frame, title, 28, first=True, bold=True)

    if position:
        marker = _textbox(slide, w - _MARGIN - 1.6, 0.62, 1.6, 0.35)
        _line(marker, position, 11, first=True, color=_MUTED, align=PP_ALIGN.RIGHT)

    _rect(slide, _MARGIN, 1.45, content_w, 0.02, _BORDER)


def _add_title_slide(presentation, result: GenerateResponse) -> None:
    slide = _blank(presentation)
    w, h = _geom(presentation)
    _rect(slide, 0, 0, 0.22, h, _ACCENT)

    request = result.request
    frame = _textbox(slide, 1.1, 2.0, w - 2.2, 1.7)
    _line(frame, result.slide_deck.title or "발표자료", 36, first=True, bold=True)

    meta = " · ".join(
        [
            labels.AUDIENCE_LABELS[request.audience],
            labels.PURPOSE_LABELS[request.purpose],
            f"{request.duration_minutes}분",
            labels.STYLE_LABELS[request.style],
            f"슬라이드 {len(result.slide_deck.slides)}장",
        ]
    )
    meta_frame = _textbox(slide, 1.1, h - 3.65, w - 2.2, 0.4)
    _line(meta_frame, meta, 15, first=True, color=_MUTED)

    notices: list[tuple[str, str]] = []
    if request.audience.value == "customer":
        notices.append(("공개 전 검토 필요", "고객용 자료입니다. 내부 정보가 없는지 확인하세요."))
    if result.verification_report is not None:
        report = result.verification_report
        notices.append(
            (
                f"원문 대비 검증: {labels.STATUS_LABELS[report.status]}",
                report.summary,
            )
        )
    if result.meta.fallback_used:
        notices.append(("기본 분석 결과로 대체됨", result.meta.fallback_reason))

    top = h - 3.05
    for label, detail in notices:
        _rect(slide, 1.1, top, 0.06, 0.5, _ACCENT)
        badge = _textbox(slide, 1.3, top + 0.02, w - 2.4, 0.5)
        _line(badge, label, 13, first=True, bold=True, color=_WARN, space_after=2)
        if detail:
            _line(badge, detail, 11, first=False, color=_MUTED)
        top += 0.62

    footer = _textbox(slide, 1.1, h - 1.0, w - 2.2, 0.5)
    _line(footer, _DISCLAIMER, 11, first=True, color=_MUTED)


def _refs_label(result: GenerateResponse, refs: list[str]) -> str:
    """근거를 파일에 적는 말. `chunk-02` 는 내부 식별자라 원문 쪽으로 바꿔 부른다.

    받아 본 사람이 원문에서 그 자리를 펴 볼 수 있어야 근거 표기가 제 역할을 한다.
    쪽을 못 찾은 근거는 지우지 않고 식별자 그대로 남긴다 — 추적을 끊는 것보다 낫다.
    화면(`frontend/components/EvidenceRef.tsx`)과 Markdown 도 같은 규칙을 쓴다.
    """
    pages = {item.id: item.page for item in result.source_analysis.source_evidence}
    labels: list[str] = []
    for ref in refs:
        label = f"{pages[ref]}쪽" if ref in pages else ref
        if label not in labels:
            labels.append(label)
    return ", ".join(labels) or "없음"


def _add_content_slide(
    presentation, slide_data: Slide, index: int, total: int, notes: str, refs: str
) -> None:
    slide = _blank(presentation)
    w, h = _geom(presentation)
    content_w = w - _MARGIN * 2
    _slide_header(slide, slide_data.title or f"슬라이드 {index}", f"{index:02d} / {total:02d}", w)

    if slide_data.takeaway:
        _rect(slide, _MARGIN, 1.7, content_w, 0.78, _ACCENT_SOFT)
        _rect(slide, _MARGIN, 1.7, 0.07, 0.78, _ACCENT)
        frame = _textbox(slide, _MARGIN + 0.25, 1.87, content_w - 0.5, 0.5)
        _line(frame, slide_data.takeaway, 15, first=True, bold=True)

    body_top = 2.75 if slide_data.takeaway else 1.75
    body_w = content_w - 0.3
    body_h = h - 1.35 - body_top
    frame = _textbox(slide, _MARGIN + 0.15, body_top, body_w, body_h)
    bullets = [f"·  {bullet}" for bullet in slide_data.bullets]
    # 바닥 크기로도 넘치면 그대로 놓는다 — 새로 그리는 경로에는 물러설 자리가 없다.
    size = _fit_body_size(bullets, body_w, body_h, space_after=10) or 9
    for order, bullet in enumerate(bullets):
        paragraph = _line(frame, bullet, size, first=order == 0, space_after=10)
        paragraph.line_spacing = 1.25

    _rect(slide, _MARGIN, h - 1.2, content_w, 0.02, _BORDER)
    footer = _textbox(slide, _MARGIN, h - 1.05, content_w, 0.7)
    _line(
        footer,
        f"추천 시각자료: {slide_data.visual_suggestion or '없음'}",
        10.5,
        first=True,
        color=_MUTED,
        space_after=3,
    )
    _line(
        footer,
        f"원문 근거: {refs}",
        10.5,
        first=False,
        color=_MUTED,
    )

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_qa_slides(presentation, result: GenerateResponse) -> None:
    qa = result.presentation_support.qa
    if not qa:
        return

    per_slide = 3
    pages = [qa[start : start + per_slide] for start in range(0, len(qa), per_slide)]
    for page_index, page in enumerate(pages, start=1):
        slide = _blank(presentation)
        w, h = _geom(presentation)
        position = f"{page_index} / {len(pages)}" if len(pages) > 1 else ""
        _slide_header(slide, "부록 · 예상 질문과 답변", position, w)

        frame = _textbox(slide, _MARGIN, 1.75, w - _MARGIN * 2, h - 3.1)
        drop = _shrink(
            sum(len(item.question) + len(item.answer) for item in page), soft=380, hard=560
        )
        first = True
        for item in page:
            _line(frame, f"Q. {item.question}", 15 - drop, first=first, bold=True, space_after=4)
            first = False
            _line(frame, f"A. {item.answer}", 13 - drop, first=False, space_after=3)
            _line(
                frame,
                f"원문 근거: {_refs_label(result, item.source_refs)}",
                10.5 - drop / 2,
                first=False,
                color=_MUTED,
                space_after=14,
            )


def _add_verification_slide(presentation, result: GenerateResponse) -> None:
    report = result.verification_report
    if report is None:
        return

    slide = _blank(presentation)
    w, h = _geom(presentation)
    content_w = w - _MARGIN * 2
    _slide_header(slide, "부록 · 원문 대비 검증", "", w)

    _rect(slide, _MARGIN, 1.7, content_w, 0.85, _SURFACE_MUTED)
    head = _textbox(slide, _MARGIN + 0.25, 1.87, content_w - 0.5, 0.6)
    _line(
        head,
        f"상태: {labels.STATUS_LABELS[report.status]} · 슬라이드 {report.checked_slides}장 대조",
        14,
        first=True,
        bold=True,
        space_after=3,
    )
    _line(head, report.summary, 11, first=False, color=_MUTED)

    frame = _textbox(slide, _MARGIN + 0.15, 2.8, content_w - 0.3, h - 4.2)
    if not report.items:
        _line(frame, "·  원문과 어긋나는 내용을 찾지 못했습니다.", 14, first=True)
    else:
        shown = report.items[:6]
        drop = _shrink(
            sum(len(item.message) + len(item.suggested_fix) for item in shown),
            soft=420,
            hard=620,
        )
        for order, item in enumerate(shown):
            head_text = (
                f"·  [{labels.SEVERITY_LABELS[item.severity]}] "
                f"{labels.ISSUE_TYPE_LABELS[item.type]}"
            )
            if item.slide_id:
                head_text += f" ({item.slide_id})"
            _line(frame, f"{head_text}: {item.message}", 13 - drop, first=order == 0, space_after=3)
            _line(
                frame,
                f"    수정 제안: {item.suggested_fix}",
                11 - drop / 2,
                first=False,
                color=_MUTED,
                space_after=10,
            )
        if len(report.items) > len(shown):
            _line(
                frame,
                f"·  이 밖에 {len(report.items) - len(shown)}건이 더 있습니다. 화면의 검증 탭에서 확인하세요.",
                11,
                first=False,
                color=_MUTED,
            )

    _rect(slide, _MARGIN, h - 1.2, content_w, 0.02, _BORDER)
    footer = _textbox(slide, _MARGIN, h - 1.05, content_w, 0.6)
    _line(
        footer,
        "이 검증은 확인이 필요한 지점을 알려 줄 뿐 사람 검토를 대체하지 않습니다.",
        10.5,
        first=True,
        color=_MUTED,
    )


def _speaker_notes(result: GenerateResponse, slide_data: Slide) -> str:
    """발표자 노트. 모듈 D 의 스크립트가 있으면 그것을 쓴다."""
    script = next(
        (item for item in result.presentation_support.scripts if item.slide_id == slide_data.id),
        None,
    )
    if script is None:
        return slide_data.speaker_notes

    parts = [f"[약 {script.duration_seconds}초]", script.script]
    if script.must_say:
        parts.append(f"꼭 말할 것: {script.must_say}")
    if slide_data.source_refs:
        parts.append(f"원문 근거: {_refs_label(result, slide_data.source_refs)}")
    return "\n\n".join(part for part in parts if part)


# --------------------------------------------------------------------------
# 원본 PPTX 위에 얹기 — 이미지·표·서식을 지키는 경로
# --------------------------------------------------------------------------


def _text_shapes(slide) -> list:
    """글을 갈아 끼울 수 있는 도형들.

    표(GraphicFrame)와 그림에는 text_frame 이 없어 저절로 빠진다 — 이 경로는 원본의 표·이미지를
    건드리지 않는다. 쪽번호·날짜·바닥글 placeholder 는 원본이 설계한 머리글/꼬리글이라 내용이
    아니므로 남긴다.
    """
    chrome = {
        PP_PLACEHOLDER.SLIDE_NUMBER,
        PP_PLACEHOLDER.DATE,
        PP_PLACEHOLDER.FOOTER,
    }
    shapes = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if (shape.width or 0) * (shape.height or 0) <= 0:
            continue
        if shape.is_placeholder and shape.placeholder_format.type in chrome:
            continue
        shapes.append(shape)
    return shapes


def _title_shape(slide):
    """제목을 넣을 도형.

    `shapes.title` 이 없는 원본이 흔하다 (빈 레이아웃 위에 텍스트 상자로 직접 만든 자료).
    그때 None 을 돌려주면 생성된 제목이 어디에도 들어가지 않고 원본 제목이 그대로 남으므로,
    맨 위 텍스트 상자를 제목 자리로 본다. 본문이 갈 곳은 남겨 둬야 하니 상자가 하나뿐이거나
    맨 위 상자가 곧 가장 큰 상자면 제목으로 쓰지 않는다.
    """
    try:
        title = slide.shapes.title
    except (AttributeError, ValueError):  # 제목 자리가 없는 레이아웃
        title = None
    if title is not None:
        return title

    shapes = _text_shapes(slide)
    if len(shapes) < 2:
        return None
    topmost = min(shapes, key=lambda shape: (shape.top or 0))
    largest = max(shapes, key=lambda shape: (shape.width or 0) * (shape.height or 0))
    return None if topmost.shape_id == largest.shape_id else topmost


def _body_shape(slide, title):
    """본문 글을 갈아 끼울 도형. 제목을 뺀 텍스트 도형 중 가장 넓은 것.

    제목은 `shape_id` 로 거른다 — `shapes.title` 은 호출할 때마다 새 proxy 를 만들어
    `is` 비교가 항상 참이 되고, 그러면 제목 상자에 본문이 들어간다.
    """
    title_id = title.shape_id if title is not None else None
    # ponytail: 면적 최댓값 휴리스틱. 원본에 더 큰 장식용 텍스트 상자가 있으면 그쪽을 고른다.
    candidates = [shape for shape in _text_shapes(slide) if shape.shape_id != title_id]
    if not candidates:
        return None
    return max(candidates, key=lambda shape: shape.width * shape.height)


def _first_run_rPr(frame):
    """원본 첫 run 의 글꼴 서식. 글만 바꾸고 서식은 물려주기 위해 복사해 둔다."""
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            return copy.deepcopy(run._r.get_or_add_rPr())  # noqa: SLF001
    return None


def _fit_size(
    base: float | None, lines: list[str], shape=None, *, assumed: float = 18
) -> float | None:
    """이 도형 안에 들어가는 글자 크기. 원본 크기보다 키우지 않고, 넘칠 때만 줄인다.

    예전에는 글자 수만 보고 정해서 좁은 상자에서는 그대로 넘쳤다 — 5.8인치 상자든 12인치
    상자든 같은 크기를 내주니 당연한 결과다. 문장을 자르지 않는 것이 원칙이라
    (docs/04-slide-planner.md) 줄일 수 있는 것은 글자 크기뿐이다.

    `base` 가 None 이면 원본이 크기를 레이아웃에서 물려받는 경우다. 물려받은 값이 얼마인지
    python-pptx 로는 알 수 없어, 넘치지 않을 크기를 명시한다 — 형식을 조금 덮어쓰더라도
    글이 상자 밖으로 나가는 것보다 낫다.
    """
    if shape is None:  # 도형을 모르면 예전처럼 글자 수로만 어림잡는다
        total = sum(len(line) for line in lines)
        if base is None:
            return 14 if total > 220 or len(lines) > 5 else None
        return max(base - 5, 10) if total > 360 else base

    # 텍스트 프레임의 기본 안쪽 여백(좌우 0.1in, 상하 0.05in)을 뺀 실제 글 자리
    width = (shape.width or 0) / _EMU_PER_INCH - 0.2
    height = (shape.height or 0) / _EMU_PER_INCH - 0.1
    if width <= 0.5 or height <= 0.2:
        return base

    size = base if base else assumed
    while size > 9 and _text_height(lines, size, width, space_after=4) > height:
        size -= 1

    if base is not None:
        return min(size, base)
    # 물려받은 크기로도 충분한지 알 수 없으므로, 줄일 필요가 없을 때만 원본에 맡긴다.
    return None if size >= assumed else size


def _replace_lines(
    frame, lines: list[str], *, bold_first: bool = False, shape=None, assumed: float = 18
) -> None:
    """텍스트 프레임의 글만 바꾼다. 첫 run 의 글꼴과 첫 문단의 글머리 서식을 물려준다."""
    rPr = _first_run_rPr(frame)
    raw_size = rPr.get("sz") if rPr is not None else None
    size = _fit_size(
        int(raw_size) / 100 if raw_size else None, lines, shape, assumed=assumed
    )

    frame.word_wrap = True
    frame.clear()
    pPr = copy.deepcopy(frame.paragraphs[0]._p.find(qn("a:pPr")))  # noqa: SLF001

    for order, line in enumerate(lines):
        if order == 0:
            paragraph = frame.paragraphs[0]
        else:
            paragraph = frame.add_paragraph()
            if pPr is not None:
                paragraph._p.insert(0, copy.deepcopy(pPr))  # noqa: SLF001
        run = paragraph.add_run()
        run.text = line
        if rPr is not None:
            run._r.replace(run._r.get_or_add_rPr(), copy.deepcopy(rPr))  # noqa: SLF001
        if size is not None:
            run.font.size = Pt(size)
        if bold_first and order == 0:
            run.font.bold = True


def _free_band(slide, height: float) -> tuple[float, float] | None:
    """도형이 하나도 없는 가장 넓은 세로 구간 (top, height).

    표나 그림만 있는 원본 슬라이드에 글을 얹을 자리를 찾을 때 쓴다. 겹치면 원본이 가려지므로
    빈 구간이 좁으면 포기한다 (호출부가 생성 레이아웃으로 새 슬라이드를 만든다).
    """
    spans = sorted(
        (
            (shape.top or 0) / _EMU_PER_INCH,
            ((shape.top or 0) + (shape.height or 0)) / _EMU_PER_INCH,
        )
        for shape in slide.shapes
    )

    best = (0.0, 0.0)
    cursor = 0.5
    for start, end in [*spans, (height - 0.5, height)]:
        if start - cursor > best[1]:
            best = (cursor, start - cursor)
        cursor = max(cursor, end)
    return best if best[1] >= 1.0 else None


def _rewrite_slide(
    slide, slide_data: Slide, notes: str, refs: str, w: float, h: float
) -> bool:
    """원본 슬라이드에서 **제목 자리와 본문 자리의 글만** 청중용 내용으로 바꾼다.

    갈아 끼우지 않은 텍스트 상자는 원본 글을 그대로 둔다. 차트의 축 이름·수치 라벨, 범례,
    각주, 좌우 2단 중 한쪽 같은 것들이 여기 해당한다. 한때는 이것들을 비웠는데, 원본이
    도표 위주면 슬라이드에서 숫자와 라벨이 통째로 사라져 발표에 쓸 수 없는 파일이 나왔다.
    비우지 않으므로 검증을 거치지 않은 원문이 완성본에 남는다 — 이 경로가 지키는 것은
    "원본 형식 보존"이고, 원문 대비 검증은 화면의 검증 탭이 계속 책임진다.

    이미지·표·배경은 애초에 텍스트 프레임이 없어 손대지 않는다.
    글을 놓을 자리를 못 찾으면 False 를 돌려준다.
    """
    title = _title_shape(slide)

    if title is not None and slide_data.title:
        _replace_lines(title.text_frame, [slide_data.title], shape=title, assumed=28)

    lines = ([slide_data.takeaway] if slide_data.takeaway else []) + list(slide_data.bullets)
    body = _body_shape(slide, title)

    if lines and body is not None:
        _replace_lines(
            body.text_frame, lines, bold_first=bool(slide_data.takeaway), shape=body
        )

    if lines and body is None:
        band = _free_band(slide, h)
        if band is None:
            return False
        top, band_height = band
        band_w = w - _MARGIN * 2
        band_h = band_height - 0.2
        bullets = [f"·  {line}" for line in lines]
        # 빈 구간에 안 들어가면 얹지 않는다. 넘친 글은 원본 그림·표와 아래 근거 줄을 덮는다.
        size = _fit_body_size(bullets, band_w, band_h, space_after=8)
        if size is None:
            return False
        frame = _textbox(slide, _MARGIN, top + 0.1, band_w, band_h)
        for order, bullet in enumerate(bullets):
            _line(
                frame,
                bullet,
                size,
                first=order == 0,
                bold=order == 0 and bool(slide_data.takeaway),
                space_after=8,
            )

    footer = _textbox(slide, _MARGIN, h - 0.42, w - _MARGIN * 2, 0.3)
    _line(
        footer,
        f"원문 근거: {refs}",
        9,
        first=True,
        color=_MUTED,
    )

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return True


def _source_slide_index(
    slide_data: Slide, pages: dict[str, int], count: int, used: set[int]
) -> int | None:
    """생성 슬라이드가 어느 원본 슬라이드에서 왔는지. PPTX 는 페이지 번호 = 슬라이드 번호다.

    0번(표지)은 후보에서 뺀다. 표지는 원본 그대로 맨 앞에 남기기로 했고, 여기서 내주면
    발표의 얼굴에 본문이 얹힌다.
    """
    counts = Counter(pages[ref] for ref in slide_data.source_refs if ref in pages)
    for page, _ in counts.most_common():
        index = page - 1
        if 0 < index < count and index not in used:
            return index
    return None


def _assign_sources(
    deck: list[Slide], pages: dict[str, int], count: int
) -> list[tuple[Slide, int]]:
    """생성 슬라이드마다 얹을 원본 슬라이드를 정한다. 결과는 덱 순서 그대로다.

    근거(`source_refs`)가 가장 많이 가리키는 원본을 먼저 쓰고, 못 찾은 슬라이드는 남은 원본을
    앞에서부터 채운다. 순서대로 채우는 뒷받침이 없으면 근거가 여러 쪽에 걸친 슬라이드가 짝을
    찾지 못해 통째로 빠지고, 완성본 장수가 사용자가 요청한 장수와 어긋난다.

    덱 순서를 유지하는 것이 핵심이다 — 원본 순서로 다시 정렬하면 청중별로 다시 짠 이야기
    순서가 원본 순서로 되돌아가, 이 제품이 하는 일이 파일에서 사라진다.
    """
    used: set[int] = set()
    matched: dict[int, int] = {}  # 덱 위치 -> 원본 위치
    pending: list[int] = []

    for position, slide_data in enumerate(deck):
        index = _source_slide_index(slide_data, pages, count, used)
        if index is None:
            pending.append(position)
            continue
        used.add(index)
        matched[position] = index

    free = iter(sorted(set(range(1, count)) - used))
    for position in pending:
        index = next(free, None)
        if index is None:
            break
        matched[position] = index

    return [(deck[position], matched[position]) for position in sorted(matched)]


def _reorder(presentation, slide_ids, order: list) -> None:
    """최종 순서대로 다시 매단다. 순서에 없는 원본 슬라이드는 관계까지 끊어 뺀다."""
    keep = {id(element) for element in order}
    for element in list(slide_ids):
        slide_ids.remove(element)
        if id(element) not in keep:
            presentation.part.drop_rel(element.rId)
    for element in order:
        slide_ids.append(element)


def _build_on_template(result: GenerateResponse, template: bytes) -> bytes:
    """원본 PPTX 를 열어 그 슬라이드 위에 결과를 얹는다.

    완성본은 원본 슬라이드로만 이루어지고, 장수는 원본을 넘지 않는다. 새로 그린 표지도,
    예상 Q&A·검증 부록도 붙이지 않는다 — 10장을 넣고 10장을 요청한 사람에게 14장을 주면
    그 파일은 발표에 그대로 쓸 수 없다. Q&A 와 검증은 화면과 Markdown·JSON 다운로드에 있다.

    표지(원본 1장)는 글까지 원본 그대로 맨 앞에 남는다.
    """
    presentation = Presentation(io.BytesIO(template))
    w, h = _geom(presentation)

    originals = list(presentation.slides)
    if not originals:
        raise ValueError("원본 PPTX 에 슬라이드가 없습니다")

    slide_ids = presentation.slides._sldIdLst  # noqa: SLF001 - 순서 변경은 공개 API 가 없다
    original_ids = list(slide_ids)
    pages = {item.id: item.page for item in result.source_analysis.source_evidence}

    # 표지는 손대지 않는다. 순서를 여기서부터 쌓으므로 항상 맨 앞이다.
    order: list = [original_ids[0]]

    deck = result.slide_deck.slides
    capacity = len(originals) - 1  # 표지를 뺀 자리
    if len(deck) > capacity:
        _log.info(
            "원본이 %d장이라 생성 슬라이드 %d장 중 %d장만 얹습니다",
            len(originals),
            len(deck),
            capacity,
        )

    for slide_data, source in _assign_sources(deck[:capacity], pages, len(originals)):
        notes = _speaker_notes(result, slide_data)
        refs = _refs_label(result, slide_data.source_refs)
        if not _rewrite_slide(originals[source], slide_data, notes, refs, w, h):
            # 글을 놓을 자리가 없는 원본이다. 새 슬라이드를 만들면 장수가 어긋나므로
            # 원본을 그대로 둔다.
            _log.info("원본 %d번 슬라이드에 글 놓을 자리가 없어 원본 그대로 둡니다", source + 1)
        order.append(original_ids[source])

    _reorder(presentation, slide_ids, order)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


# --------------------------------------------------------------------------
# 진입점
# --------------------------------------------------------------------------


def build_pptx(result: GenerateResponse, template: bytes | None = None) -> bytes:
    """생성 결과를 PPTX 바이트로 만든다.

    표지 → 슬라이드 → 부록(예상 Q&A, 검증) 순이며, 발표자 노트에는 모듈 D 의 스크립트가 들어간다.

    `template` 은 업로드한 원본 PPTX 다. 있으면 그 파일 위에 결과를 얹어 원본의 이미지·표·
    배경·글꼴을 지킨다. 원본이 특이해 실패하면 빈 프레젠테이션에 새로 그린다 — 다운로드가
    통째로 실패하는 것보다 낫다.
    """
    if not PPTX_AVAILABLE:
        raise PptxUnavailableError(
            "PPTX 변환 모듈(python-pptx)이 설치되어 있지 않습니다. "
            "Markdown 또는 JSON 다운로드를 이용하세요."
        )

    if template:
        try:
            return _build_on_template(result, template)
        except Exception:  # noqa: BLE001 - 원본 형식은 통제 밖이다
            _log.warning("원본 PPTX 위에 얹지 못해 새 덱으로 만듭니다", exc_info=True)

    presentation = Presentation()
    presentation.slide_width = Inches(_SLIDE_W)
    presentation.slide_height = Inches(_SLIDE_H)

    _add_title_slide(presentation, result)

    slides = result.slide_deck.slides
    for index, slide_data in enumerate(slides, start=1):
        _add_content_slide(
            presentation,
            slide_data,
            index,
            len(slides),
            _speaker_notes(result, slide_data),
            _refs_label(result, slide_data.source_refs),
        )

    _add_qa_slides(presentation, result)
    _add_verification_slide(presentation, result)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def filename_for(result: GenerateResponse) -> str:
    """다운로드 파일명. 파일명에 못 쓰는 문자만 걸러 낸다."""
    title = result.slide_deck.title or result.presentation_id
    cleaned = "".join(" " if char in '\\/:*?"<>|\n\r\t' else char for char in title).strip()
    cleaned = " ".join(cleaned.split())
    return f"{cleaned or result.presentation_id}.pptx"
