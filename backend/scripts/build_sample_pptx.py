"""샘플 기술문서를 PPTX 입력 fixture 로 만든다 (docs/02-document-analysis.md).

`fixtures/sample_document.txt` 와 같은 내용을 발표자료 형태로 옮긴 것이다.
숫자·용어를 그대로 유지해야 파싱 테스트의 근거 probe 가 성립한다.
표와 발표자 노트를 일부러 섞어 두어 파서가 두 경로를 모두 타게 한다.

사용:
    backend\\.venv\\Scripts\\python.exe backend\\scripts\\build_sample_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

BACKEND = Path(__file__).resolve().parents[1]
OUTPUT = BACKEND / "fixtures" / "sample_document.pptx"

TITLE = "DocClassifier v2.1 기술 개요서"
SUBTITLE = "작성: 정보시스템팀 문서플랫폼파트\n문서 등급: 사내 공유\n최종 수정: 2026년 7월 14일"

# (제목, 본문 bullet, 발표자 노트)
SECTIONS: list[tuple[str, list[str], str]] = [
    (
        "1. 배경",
        [
            "사내에서 하루에 생성되는 업무 문서는 약 8,200건이며, 이 중 62%가 담당 부서 지정 없이 공용 저장소에 업로드된다.",
            "담당자가 문서를 열어 분류하는 데 건당 평균 3분 20초가 소요되며, 이 작업만으로 월 약 340시간의 인력이 소모된다.",
            "DocClassifier는 업로드된 문서를 사전 정의된 24개 업무 카테고리 중 하나로 분류해 담당 부서 큐로 전달한다.",
            "v1.0은 2025년 11월에 배포되었고, v2.1은 분류 정확도와 처리량을 개선한 후속 버전이다.",
        ],
        "",
    ),
    (
        "2. 시스템 구성",
        [
            "추출 단계: 업로드된 파일에서 텍스트를 추출한다. PDF, DOCX, TXT, HWP 네 가지 형식을 지원한다. "
            "텍스트 레이어가 없는 스캔 이미지 PDF는 OCR 전처리 모듈로 우회한다.",
            "표현 단계: 추출한 텍스트를 임베딩(embedding, 문장을 숫자 벡터로 바꾼 표현)으로 변환한다. "
            "문서 전체가 아니라 앞부분 2,000자만 사용한다.",
            "판정 단계: 임베딩 기반 분류기와 규칙 기반 분류기를 함께 쓰는 앙상블 구조이며, "
            "두 분류기의 판단이 엇갈리면 규칙 기반 결과를 우선한다.",
        ],
        "앞부분 2,000자만 쓰는 이유는 실험 결과 2,000자를 넘겨도 정확도 개선폭이 0.4%p 미만이었기 때문이다.",
    ),
    (
        "4. 운영 제약 사항",
        [
            "개인정보가 포함된 문서는 마스킹 처리를 거친 뒤에만 분류기에 입력해야 한다. "
            "마스킹 모듈을 거치지 않은 원문을 직접 전달하는 것은 금지되어 있다.",
            "한 번에 3만 건을 초과하는 대량 분류 요청은 실시간 API가 아니라 야간 배치로 분리해야 한다.",
            "신규 카테고리를 추가하려면 해당 카테고리 학습 데이터가 최소 200건 이상 필요하다.",
            "분류 신뢰도가 0.6 미만인 문서는 자동 분류하지 않고 담당자 검토 큐로 보낸다. 전체 문서의 약 7%가 이 경로로 처리된다.",
        ],
        "",
    ),
    (
        "5. 도입 효과",
        [
            "문서플랫폼파트가 3개 부서를 대상으로 8주간 운영한 결과, 문서 분류 소요 시간이 건당 평균 3분 20초에서 25초로 줄었다.",
            "같은 기간 오분류로 인한 문서 재이관 건수는 주당 평균 41건에서 9건으로 감소했다.",
        ],
        "25초는 담당자 검토 큐로 넘어가는 7%를 포함한 수치다.",
    ),
    (
        "6. 향후 계획",
        [
            "2026년 4분기에 사내 K-Drive와 연동해 전사 확대를 검토하고 있다.",
            "연동 시 업로드 시점에 분류가 완료되므로 별도 업로드 단계가 필요 없어진다.",
            "다만 K-Drive 연동은 정보보호팀의 보안성 검토 승인을 전제로 하며, 승인 일정은 아직 확정되지 않았다.",
        ],
        "",
    ),
]

PERFORMANCE_ROWS = [
    ("지표", "v2.1", "비고"),
    ("분류 정확도", "94.2%", "v1.0은 87.6%"),
    ("매크로 F1 점수", "0.91", "정밀도와 재현율의 조화평균"),
    ("평균 응답 시간", "210ms", "텍스트 레이어가 있는 문서 기준"),
    ("최대 처리량", "초당 45건", "사내 테스트셋 12,400건 기준"),
]

PERFORMANCE_NOTE = (
    "정확도 94.2%는 텍스트 레이어가 있는 문서에 한한 수치다. "
    "스캔 이미지 PDF는 OCR 품질에 따라 정확도가 76%까지 떨어지며, 평균 응답 시간도 1.8초로 증가한다."
)


def _set_notes(slide, text: str) -> None:
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _add_title_slide(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = SUBTITLE


def _add_bullet_slide(presentation: Presentation, title: str, bullets: list[str], notes: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title

    frame = slide.placeholders[1].text_frame
    frame.text = bullets[0]
    for bullet in bullets[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = bullet
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    _set_notes(slide, notes)


def _add_table_slide(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "3. 성능"

    rows, cols = len(PERFORMANCE_ROWS), len(PERFORMANCE_ROWS[0])
    left, top = presentation.slide_width // 10, presentation.slide_height // 4
    width = presentation.slide_width - 2 * left
    height = presentation.slide_height // 3
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for r, row in enumerate(PERFORMANCE_ROWS):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

    _set_notes(slide, PERFORMANCE_NOTE)


def main() -> int:
    presentation = Presentation()

    _add_title_slide(presentation)
    for title, bullets, notes in SECTIONS[:2]:
        _add_bullet_slide(presentation, title, bullets, notes)
    _add_table_slide(presentation)
    for title, bullets, notes in SECTIONS[2:]:
        _add_bullet_slide(presentation, title, bullets, notes)

    presentation.save(OUTPUT)
    print(f"생성: fixtures/{OUTPUT.name} (슬라이드 {len(presentation.slides._sldIdLst)}장)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
