"""PPTX export (docs/04-slide-planner.md).

export 는 새 문장을 만들지 않는다. 화면에 보이는 내용이 파일 안에도 그대로 있어야 한다.
"""

from __future__ import annotations

import base64
import io
import zipfile
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from app.api.presentations import PPTX_MEDIA_TYPE
from app.main import app
from app.models.contracts import GenerateResponse, SourceEvidence
from app.services import export_pptx
from app.services.pipeline import build_document, generate
from app.services.store import store


@pytest.fixture()
def result(sample_text: str, customer_request) -> GenerateResponse:
    document = build_document("sample_document.txt", sample_text.encode("utf-8"))
    return generate(document, customer_request)


def _open(data: bytes) -> Presentation:
    return Presentation(io.BytesIO(data))


def _all_text(presentation: Presentation) -> str:
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def test_deck_has_cover_slides_and_appendix(result: GenerateResponse) -> None:
    presentation = _open(export_pptx.build_pptx(result))

    # 표지 1장 + 본문 + 예상 Q&A 부록 + 검증 부록
    assert len(presentation.slides) > len(result.slide_deck.slides) + 1

    text = _all_text(presentation)
    assert result.slide_deck.title in text
    for slide in result.slide_deck.slides:
        assert slide.takeaway in text


def test_body_text_is_sized_to_fit_its_band() -> None:
    """넘친 본문이 아래 '원문 근거' 줄을 덮으면 안 된다.

    문장은 자르지 않는 것이 원칙이라 줄일 수 있는 건 글자 크기뿐이고, 그것으로도 안 되면
    호출부가 그 슬라이드를 포기한다(`_rewrite_slide` 가 False 를 돌려준다).
    """
    lines = ["·  " + "원문에서 그대로 가져온 긴 문장입니다. " * 6] * 4

    fits = export_pptx._fit_body_size(lines, 12.0, 4.0, space_after=8)
    assert fits is not None
    assert export_pptx._text_height(lines, fits, 12.0, space_after=8) <= 4.0

    # 자리가 없으면 억지로 밀어 넣지 않고 None 을 돌려준다.
    assert export_pptx._fit_body_size(lines, 12.0, 0.6, space_after=8) is None


def test_bullets_are_not_truncated(result: GenerateResponse) -> None:
    """문장을 자르면 깨진 어미가 남는다. 넘칠 때는 글자 크기를 줄인다."""
    text = _all_text(_open(export_pptx.build_pptx(result)))
    for slide in result.slide_deck.slides:
        for bullet in slide.bullets:
            assert bullet in text


def test_source_refs_are_printed_on_every_slide(result: GenerateResponse) -> None:
    """데모 성공 기준 4번 — 근거는 파일에서도 따라갈 수 있어야 한다.

    파일을 받아 보는 사람에게 `chunk-02` 는 아무 뜻이 없으므로 원문 쪽으로 적는다.
    """
    presentation = _open(export_pptx.build_pptx(result))
    # 표지 다음부터가 본문 슬라이드다
    body = list(presentation.slides)[1 : 1 + len(result.slide_deck.slides)]
    pages = {item.id: item.page for item in result.source_analysis.source_evidence}

    for slide_data, slide in zip(result.slide_deck.slides, body):
        text = "\n".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        assert "원문 근거:" in text
        assert "chunk-" not in text
        for ref in slide_data.source_refs:
            assert f"{pages[ref]}쪽" in text


def test_speaker_notes_carry_the_script(result: GenerateResponse) -> None:
    presentation = _open(export_pptx.build_pptx(result))
    body = list(presentation.slides)[1 : 1 + len(result.slide_deck.slides)]

    scripts = {item.slide_id: item for item in result.presentation_support.scripts}
    for slide_data, slide in zip(result.slide_deck.slides, body):
        notes = slide.notes_slide.notes_text_frame.text
        script = scripts.get(slide_data.id)
        assert script is not None
        assert script.script in notes


def test_customer_deck_warns_before_publishing(result: GenerateResponse) -> None:
    """고객 청중이면 '공개 전 검토 필요' 를 색이 아니라 글자로 남긴다."""
    assert result.request.audience.value == "customer"
    assert "공개 전 검토 필요" in _all_text(_open(export_pptx.build_pptx(result)))


def test_verification_status_uses_text_label(result: GenerateResponse) -> None:
    from app.services import labels

    report = result.verification_report
    assert report is not None
    text = _all_text(_open(export_pptx.build_pptx(result)))
    assert labels.STATUS_LABELS[report.status] in text


def test_filename_drops_characters_windows_rejects(result: GenerateResponse) -> None:
    result.slide_deck.title = 'a/b:c*d?e"f<g>h|i'
    assert export_pptx.filename_for(result) == "a b c d e f g h i.pptx"


def test_filename_falls_back_to_presentation_id(result: GenerateResponse) -> None:
    result.slide_deck.title = ""
    assert export_pptx.filename_for(result) == f"{result.presentation_id}.pptx"


# --- API ------------------------------------------------------------------


@pytest.fixture()
def client() -> TestClient:
    store.clear()
    return TestClient(app)


def _generate(client: TestClient, sample_text: str) -> str:
    uploaded = client.post(
        "/api/documents",
        files={"file": ("sample_document.txt", sample_text.encode("utf-8"), "text/plain")},
    ).json()
    generated = client.post(
        "/api/presentations/generate",
        json={
            "document_id": uploaded["document"]["document_id"],
            "request": {
                "audience": "executive",
                "purpose": "internal_report",
                "duration_minutes": 5,
                "keywords": [],
                "style": "concise",
                "preserve_original_terms": True,
                "slide_count": 5,
            },
        },
    )
    assert generated.status_code == 200, generated.text
    return generated.json()["presentation_id"]


def test_export_endpoint_returns_a_real_pptx(client: TestClient, sample_text: str) -> None:
    presentation_id = _generate(client, sample_text)

    response = client.get(f"/api/presentations/{presentation_id}/export/pptx")
    assert response.status_code == 200, response.text
    assert response.headers["content-type"] == PPTX_MEDIA_TYPE

    # PPTX 는 zip 이다. 실제로 열리는지까지 확인한다.
    assert zipfile.is_zipfile(io.BytesIO(response.content))
    assert len(_open(response.content).slides) > 1


def test_export_sets_a_downloadable_korean_filename(
    client: TestClient, sample_text: str
) -> None:
    presentation_id = _generate(client, sample_text)
    response = client.get(f"/api/presentations/{presentation_id}/export/pptx")

    disposition = response.headers["content-disposition"]
    assert disposition.startswith("attachment;")
    # 한국어 파일명은 ASCII 로 못 쓴다. 대체 이름과 RFC 5987 이름이 함께 있어야 한다.
    assert f'filename="{presentation_id}.pptx"' in disposition
    assert "filename*=UTF-8''" in disposition


def test_export_unknown_presentation_returns_404(client: TestClient) -> None:
    response = client.get("/api/presentations/pres-missing/export/pptx")
    assert response.status_code == 404
    assert "발표 결과를 찾을 수 없습니다" in response.json()["detail"]


def test_export_without_python_pptx_explains_in_korean(
    client: TestClient, sample_text: str, monkeypatch: pytest.MonkeyPatch
) -> None:
    """의존성이 빠져도 앱은 살아 있고, 다른 다운로드로 안내한다."""
    presentation_id = _generate(client, sample_text)
    monkeypatch.setattr(export_pptx, "PPTX_AVAILABLE", False)

    response = client.get(f"/api/presentations/{presentation_id}/export/pptx")
    assert response.status_code == 503
    assert "Markdown 또는 JSON 다운로드" in response.json()["detail"]


# --- 원본 PPTX 형식 지키기 --------------------------------------------------
#
# 입력이 PPTX 면 export 는 새 파일을 그리지 않고 원본 위에 얹는다. 원본의 이미지·표는
# 텍스트 프레임이 아니므로 손대지 않는 것이 규칙이고, 이 테스트가 그 규칙을 지킨다.

_PIXEL_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _add_cover(presentation) -> None:
    """원본 표지. export 는 이 슬라이드를 글까지 그대로 맨 앞에 남겨야 한다."""
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = "원본표지제목"
    cover.placeholders[1].text_frame.text = "원본표지부제"


def _template_with_picture_and_table() -> bytes:
    """3장짜리 원본. 1장은 표지, 3장은 표·그림만 있고 본문 상자가 없다 (빈 자리 탐색 경로)."""
    presentation = Presentation()
    _add_cover(presentation)

    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "원본 제목"
    first.placeholders[1].text_frame.text = "원본 본문"

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "원본 표 슬라이드"
    second.shapes.add_table(2, 2, Inches(0.6), Inches(1.8), Inches(4.0), Inches(1.2))
    second.shapes.add_picture(
        io.BytesIO(_PIXEL_PNG), Inches(5.2), Inches(1.8), Inches(2.0), Inches(1.2)
    )

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _count_visuals(presentation: Presentation) -> tuple[int, int]:
    pictures = tables = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if getattr(shape, "has_table", False):
                tables += 1
    return pictures, tables


@pytest.fixture()
def two_slide_result(result: GenerateResponse) -> GenerateResponse:
    """생성 슬라이드 2장이 원본 2·3장에서 왔다고 두고 근거를 다시 건다.

    1장은 표지라 후보가 아니다 — 근거가 1장을 가리켜도 본문이 표지를 덮지 않는다.
    """
    result.source_analysis.source_evidence = [
        SourceEvidence(id="chunk-01", text="", page=2),
        SourceEvidence(id="chunk-02", text="", page=3),
    ]
    result.slide_deck.slides = result.slide_deck.slides[:2]
    for slide_data, ref in zip(result.slide_deck.slides, ["chunk-01", "chunk-02"]):
        slide_data.source_refs = [ref]
    return result


def test_original_pictures_and_tables_survive(two_slide_result: GenerateResponse) -> None:
    """사용자가 만들어 둔 이미지·표가 사라지지 않는다 — 이 모듈의 존재 이유."""
    template = _template_with_picture_and_table()
    assert _count_visuals(_open(template)) == (1, 1)

    exported = _open(export_pptx.build_pptx(two_slide_result, template=template))
    assert _count_visuals(exported) == (1, 1)


def test_generated_text_replaces_the_original_text(
    two_slide_result: GenerateResponse,
) -> None:
    """형식만 남기고 글은 청중용으로 바뀌어야 한다. 표만 있는 슬라이드도 마찬가지."""
    template = _template_with_picture_and_table()
    text = _all_text(_open(export_pptx.build_pptx(two_slide_result, template=template)))

    assert "원본 본문" not in text
    for slide_data in two_slide_result.slide_deck.slides:
        for bullet in slide_data.bullets:
            assert bullet in text


def test_title_only_slide_keeps_the_title_in_its_title_box(
    two_slide_result: GenerateResponse,
) -> None:
    """제목 상자에 본문 bullet 이 들어가면 안 된다.

    `shapes.title` 은 호출할 때마다 새 proxy 라 `is` 로 거르면 제목이 본문 후보가 되고,
    표만 있는 슬라이드에서는 제목 자리에 bullet 이 통째로 들어간다.
    """
    template = _template_with_picture_and_table()
    exported = _open(export_pptx.build_pptx(two_slide_result, template=template))

    slide_data = two_slide_result.slide_deck.slides[1]
    table_slide = exported.slides[2]  # 원본 표지 → 원본 2장 → 원본 3장(표 슬라이드)
    assert table_slide.shapes.title.text_frame.text == slide_data.title


def test_export_never_grows_past_the_original(two_slide_result: GenerateResponse) -> None:
    """원본 3장을 넣었으면 완성본도 3장이다.

    예전에는 새 표지 1장과 예상 Q&A·검증 부록이 붙어 요청한 장수를 넘겼다. 10장을 넣고
    10장을 요청한 사람이 14장을 받으면 그 파일은 발표에 그대로 쓸 수 없다.
    """
    template = _template_with_picture_and_table()
    source = _open(template)
    exported = _open(export_pptx.build_pptx(two_slide_result, template=template))

    assert len(exported.slides) == len(source.slides) == 3
    assert len(exported.slides) <= len(two_slide_result.slide_deck.slides) + 1

    text = _all_text(exported)
    assert "예상 질문과 답변" not in text
    assert "원문 대비 검증" not in text


def test_deck_longer_than_the_original_is_capped(two_slide_result: GenerateResponse) -> None:
    """생성 슬라이드가 원본보다 많아도 새 슬라이드를 만들지 않는다."""
    two_slide_result.slide_deck.slides = two_slide_result.slide_deck.slides * 3
    template = _template_with_picture_and_table()

    exported = _open(export_pptx.build_pptx(two_slide_result, template=template))
    assert len(exported.slides) == 3  # 표지 + 원본 2장


def test_cover_slide_is_kept_untouched(two_slide_result: GenerateResponse) -> None:
    """표지는 맨 앞에 원본 글까지 그대로 남는다."""
    exported = _open(
        export_pptx.build_pptx(two_slide_result, template=_template_with_picture_and_table())
    )

    cover = exported.slides[0]
    text = "\n".join(
        shape.text_frame.text for shape in cover.shapes if shape.has_text_frame
    )
    assert "원본표지제목" in text
    assert "원본표지부제" in text
    for slide_data in two_slide_result.slide_deck.slides:
        assert slide_data.title not in text


def test_template_keeps_its_own_slide_size(two_slide_result: GenerateResponse) -> None:
    """원본이 4:3 이면 부록 슬라이드도 4:3 이어야 한다. 16:9 로 늘리면 원본이 잘린다."""
    template = _template_with_picture_and_table()
    source = _open(template)
    exported = _open(export_pptx.build_pptx(two_slide_result, template=template))

    assert exported.slide_width == source.slide_width
    assert exported.slide_height == source.slide_height


def test_broken_template_falls_back_to_a_new_deck(result: GenerateResponse) -> None:
    """원본을 못 열어도 다운로드는 성공해야 한다."""
    exported = _open(export_pptx.build_pptx(result, template=b"not a pptx"))
    assert len(exported.slides) > 1


def test_uploaded_pptx_is_kept_for_export(customer_request) -> None:
    """plumbing: 원본 바이트를 들고 있지 않으면 얹을 것이 없다."""
    fixture = Path(__file__).resolve().parents[1] / "fixtures" / "sample_document.pptx"
    if not fixture.exists():
        pytest.skip("scripts/build_sample_pptx.py 로 fixture 를 먼저 만들어 주세요.")

    data = fixture.read_bytes()
    document = build_document("sample_document.pptx", data)
    assert document.source == data

    # TXT 입력은 얹을 원본이 없다
    assert build_document("sample.txt", b"a" * 200).source is None

    exported = _open(export_pptx.build_pptx(generate(document, customer_request), template=data))
    assert _count_visuals(exported)[1] >= 1  # 원본의 성능 표가 살아 있다


# --------------------------------------------------------------------------
# 텍스트 상자가 여럿인 원본 (빈 레이아웃 위에 직접 만든 사내 자료 모양)
# --------------------------------------------------------------------------


def _template_with_many_text_boxes() -> bytes:
    """제목 placeholder 없이 텍스트 상자만 여럿 놓인 원본.

    실제 사내 자료에 흔한 모양이다. 제목 자리를 못 찾고, 가장 큰 상자 하나만 갈아 끼우던
    시절에는 나머지 상자에 원문이 그대로 남아 새 글과 겹쳐 보였다.
    """
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    _add_cover(presentation)
    blank = presentation.slide_layouts[6]

    for number in (1, 2):
        slide = presentation.slides.add_slide(blank)

        def box(left, top, width, height, text, size):
            shape = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.text_frame.text = text
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(size)

        box(0.7, 0.5, 12.0, 1.0, f"원본제목{number}", 36)
        # 본문 자리를 넉넉하지 않게 둔다 — 글자 크기를 도형에 맞추지 않으면 그대로 넘친다.
        box(0.7, 1.8, 5.8, 2.2, f"원본왼쪽본문{number}", 18)
        box(6.8, 1.8, 5.8, 2.2, f"원본오른쪽본문{number}", 18)
        box(0.7, 6.6, 12.0, 0.5, f"원본하단주석{number}", 12)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_untouched_text_boxes_keep_their_original_text(
    two_slide_result: GenerateResponse,
) -> None:
    """갈아 끼우지 않은 상자의 원본 글은 그대로 남아야 한다.

    한때는 이 상자들을 비웠다. 그러면 원본이 도표 위주일 때 축 이름·수치 라벨·범례가 통째로
    사라져 발표에 못 쓰는 파일이 나온다. 바꾸는 것은 제목 자리와 본문 자리뿐이다.
    """
    exported = _open(
        export_pptx.build_pptx(two_slide_result, template=_template_with_many_text_boxes())
    )
    text = _all_text(exported)

    # 제목 자리(맨 위)와 본문 자리(가장 큰 상자)는 갈아 끼우고, 나머지는 손대지 않는다.
    for kept in ("원본오른쪽본문1", "원본하단주석1", "원본오른쪽본문2", "원본하단주석2"):
        assert kept in text, f"'{kept}' 가 사라졌습니다"


def test_title_lands_even_without_a_title_placeholder(
    two_slide_result: GenerateResponse,
) -> None:
    """제목 placeholder 가 없으면 맨 위 텍스트 상자가 제목 자리다.

    None 을 돌려주던 시절에는 생성된 제목이 어디에도 들어가지 않고 원본 제목이 그대로 남았다.
    """
    exported = _open(
        export_pptx.build_pptx(two_slide_result, template=_template_with_many_text_boxes())
    )
    titles = {slide_data.title for slide_data in two_slide_result.slide_deck.slides}
    text = _all_text(exported)
    for title in titles:
        assert title in text


def test_body_text_fits_inside_its_shape(two_slide_result: GenerateResponse) -> None:
    """글자 크기를 도형 크기에 맞춘다. 글자 수만 보고 정하면 좁은 상자에서 그대로 넘친다."""
    exported = _open(
        export_pptx.build_pptx(two_slide_result, template=_template_with_many_text_boxes())
    )

    checked = 0
    for slide in exported.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            lines = [line for line in shape.text_frame.text.split("\n") if line]
            sizes = [
                run.font.size.pt
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.font.size is not None
            ]
            if not lines or not sizes:
                continue
            width = shape.width / export_pptx._EMU_PER_INCH - 0.2  # noqa: SLF001
            height = shape.height / export_pptx._EMU_PER_INCH - 0.1  # noqa: SLF001
            if width <= 0.5 or height <= 0.2:
                continue
            needed = export_pptx._text_height(  # noqa: SLF001
                lines, max(sizes), width, space_after=4
            )
            assert needed <= height + 0.35, (
                f"'{lines[0][:20]}...' 가 {shape.width / export_pptx._EMU_PER_INCH:.1f}"  # noqa: SLF001
                f"x{shape.height / export_pptx._EMU_PER_INCH:.1f}in 상자를 넘칩니다"  # noqa: SLF001
            )
            checked += 1

    assert checked > 0, "검사한 텍스트 상자가 없습니다"


# --------------------------------------------------------------------------
# 머리글·차트 라벨이 있는 원본 (실무 IR 자료 모양)
# --------------------------------------------------------------------------


def _template_like_an_ir_deck() -> bytes:
    """머리글 한 줄, 그 아래 진짜 제목, 큰 숫자 라벨, 작은 차트 라벨.

    실제 IR 자료의 모양이다. "맨 위 상자 = 제목", "가장 큰 상자 = 본문" 으로 고르던 시절에는
    12pt 머리글이 발표 제목으로 바뀌고 25pt 원본 제목은 그대로 남아 한 장에 제목이 둘이 됐고,
    본문은 '92,929억원' 같은 차트 라벨을 덮어 원본 숫자를 지웠다.
    """
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    _add_cover(presentation)
    blank = presentation.slide_layouts[6]

    for number in (1, 2):
        slide = presentation.slides.add_slide(blank)

        def box(left, top, width, height, text, size):
            shape = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.text_frame.text = text
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(size)

        # 맨 위 두 줄은 머리글이다 — 이동 경로와 문서 제목.
        box(0.23, 0.08, 1.93, 0.23, f"머리글경로{number}", 12)
        box(9.74, 0.07, 3.03, 0.23, f"머리글문서명{number}", 12)
        # 진짜 제목은 그 아래에 더 큰 글자로 있다.
        box(0.24, 0.72, 2.81, 0.44, f"원본진짜제목{number}", 25)
        # 차트 가운데의 큰 숫자와 작은 축 라벨.
        box(0.17, 1.61, 2.39, 1.48, f"92,929억원 영업이익 13,655 ({number})", 45)
        box(9.32, 2.33, 0.43, 0.16, f"차트라벨{number}", 8)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_title_replaces_the_real_title_not_the_header_line(
    two_slide_result: GenerateResponse,
) -> None:
    """제목은 맨 위 머리글이 아니라 위쪽 띠에서 글자가 가장 큰 상자로 간다.

    머리글에 써 넣으면 원본 제목이 그대로 남아 한 장에 제목이 둘이 된다 — 사용자가 본
    "변경 전 텍스트가 그대로 남아있다" 가 이것이다.
    """
    exported = _open(
        export_pptx.build_pptx(two_slide_result, template=_template_like_an_ir_deck())
    )
    text = _all_text(exported)

    for number in (1, 2):
        assert f"원본진짜제목{number}" not in text, "원본 제목이 갈아 끼워지지 않았습니다"
        assert f"머리글경로{number}" in text, "머리글이 제목으로 덮였습니다"
        assert f"머리글문서명{number}" in text, "머리글이 제목으로 덮였습니다"

    for slide_data in two_slide_result.slide_deck.slides:
        assert slide_data.title in text


def test_chart_labels_are_never_used_as_the_body_box(
    two_slide_result: GenerateResponse,
) -> None:
    """차트 라벨에는 본문을 넣지 않는다. 넣으면 원본 숫자가 통째로 사라진다."""
    exported = _open(
        export_pptx.build_pptx(two_slide_result, template=_template_like_an_ir_deck())
    )
    text = _all_text(exported)

    for number in (1, 2):
        assert f"92,929억원 영업이익 13,655 ({number})" in text, "차트 숫자가 사라졌습니다"
        assert f"차트라벨{number}" in text, "차트 축 라벨이 사라졌습니다"


# --------------------------------------------------------------------------
# 짝짓기 알리기 — 화면의 "원본과 결과 비교" 가 이것만 믿는다
# --------------------------------------------------------------------------


def test_source_map_matches_where_the_slides_actually_land(
    two_slide_result: GenerateResponse,
) -> None:
    """`source_map` 이 알려 준 자리에 실제로 그 슬라이드가 있어야 한다.

    화면이 짝짓기를 다시 구현하던 시절에는 표지를 후보에서 빼는 규칙이 화면에 없어,
    원본 2장에 얹힌 슬라이드를 원본 1장(표지) 옆에 놓았다.
    """
    template = _template_with_picture_and_table()
    mapping = export_pptx.source_map(two_slide_result, template)
    exported = _open(export_pptx.build_pptx(two_slide_result, template=template))

    assert mapping["source_slides"] == len(_open(template).slides)
    assert mapping["cover_page"] == 1

    for pair in mapping["pairs"]:
        slide_data = two_slide_result.slide_deck.slides[pair["number"] - 1]
        # 표지는 후보가 아니다.
        assert pair["page"] != mapping["cover_page"]
        text = "\n".join(
            shape.text_frame.text
            for shape in exported.slides[pair["output"] - 1].shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert slide_data.title in text, (
            f"발표용 {pair['number']}장이 파일 {pair['output']}장에 있다고 했지만 없습니다"
        )


def test_source_map_without_a_template_counts_the_generated_cover(
    result: GenerateResponse,
) -> None:
    """원본이 PPTX 가 아니면 새로 그린 표지 뒤로 덱 순서 그대로다."""
    mapping = export_pptx.source_map(result, None)

    assert mapping["source_slides"] == 0
    assert mapping["cover_page"] is None
    assert [pair["output"] for pair in mapping["pairs"]] == [
        number + 1 for number in range(1, len(result.slide_deck.slides) + 1)
    ]
    assert all(pair["page"] is None for pair in mapping["pairs"])


def test_source_map_endpoint_answers_for_a_stored_presentation(
    client: TestClient, sample_text: str
) -> None:
    presentation_id = _generate(client, sample_text)

    response = client.get(f"/api/presentations/{presentation_id}/source-map")
    assert response.status_code == 200, response.text
    assert response.json()["pairs"]

    assert client.get("/api/presentations/nope/source-map").status_code == 404
