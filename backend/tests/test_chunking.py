"""모듈 A: 파싱과 청킹 (docs/02-document-analysis.md)."""

from __future__ import annotations

from pathlib import Path

import pytest

from app.models.contracts import Chunk
from app.services.chunking import build_chunks, make_chunk_id
from app.services.document_parser import DocumentError, parse_document


def test_chunk_id_format() -> None:
    assert make_chunk_id(0) == "chunk-01"
    assert make_chunk_id(11) == "chunk-12"


def test_chunks_are_created_with_pages(chunks: list[Chunk]) -> None:
    assert len(chunks) >= 3
    assert all(chunk.page >= 1 for chunk in chunks)
    assert [chunk.id for chunk in chunks] == [make_chunk_id(i) for i in range(len(chunks))]


def test_chunks_preserve_source_content(chunks: list[Chunk], sample_text: str) -> None:
    """근거로 인용할 핵심 구절이 청킹 과정에서 사라지면 안 된다."""
    joined = "\n".join(chunk.text for chunk in chunks)
    for probe in ("94.2%", "12,400건", "0.6 미만", "200건 이상", "K-Drive", "41건에서 9건"):
        assert probe in joined, f"청킹 후 '{probe}' 가 사라졌습니다"


def test_pdf_input_produces_same_evidence(sample_text: str) -> None:
    """PDF 는 이 도구의 주 입력 형식이다. 한글 추출과 페이지 정보가 유지되어야 한다.

    fixture PDF 는 sample_document.txt 를 그대로 옮긴 것이다.
    """
    pdf_path = Path(__file__).resolve().parents[1] / "fixtures" / "sample_document.pdf"
    if not pdf_path.exists():
        pytest.skip("PDF fixture 가 없습니다")

    pages = parse_document("sample_document.pdf", pdf_path.read_bytes())
    pdf_chunks = build_chunks(pages)

    assert len(pages) >= 2
    assert {chunk.page for chunk in pdf_chunks} == {page.page for page in pages}

    joined = "\n".join(chunk.text for chunk in pdf_chunks)
    for probe in ("94.2%", "12,400건", "마스킹", "3만 건", "K-Drive"):
        assert probe in joined, f"PDF 추출에서 '{probe}' 가 사라졌습니다"


def test_pptx_input_keeps_slide_pages_and_evidence() -> None:
    """PPTX 는 발표자료를 다시 만드는 입력이라 자주 들어온다.

    슬라이드 1장 = 페이지 1개이고, 표와 발표자 노트의 숫자까지 근거로 남아야 한다.
    fixture 는 scripts/build_sample_pptx.py 가 sample_document.txt 를 옮겨 만든 것이다.
    """
    pptx_path = Path(__file__).resolve().parents[1] / "fixtures" / "sample_document.pptx"
    if not pptx_path.exists():
        pytest.skip("PPTX fixture 가 없습니다")

    pages = parse_document("sample_document.pptx", pptx_path.read_bytes())
    assert [page.page for page in pages] == list(range(1, len(pages) + 1))
    assert len(pages) >= 5

    joined = "\n".join(chunk.text for chunk in build_chunks(pages))
    for probe in (
        "94.2%",  # 표 안의 숫자
        "12,400건",  # 표 안의 숫자
        "0.4%p",  # 발표자 노트 안의 숫자
        "마스킹",
        "3만 건",
        "K-Drive",
    ):
        assert probe in joined, f"PPTX 추출에서 '{probe}' 가 사라졌습니다"


def _pptx_bytes(build) -> bytes:  # noqa: ANN001 - 테스트 헬퍼
    from io import BytesIO

    from pptx import Presentation

    presentation = Presentation()
    build(presentation)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_pptx_reads_shapes_in_visual_order() -> None:
    """z-order 가 아니라 배치 순서로 읽어야 제목이 본문 앞에 온다."""
    from pptx.util import Emu

    def build(presentation) -> None:  # noqa: ANN001
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        # 아래쪽 본문을 먼저(=z-order 앞) 넣는다
        body = slide.shapes.add_textbox(Emu(914400), Emu(3657600), Emu(4572000), Emu(914400))
        body.text_frame.text = "본문 문장입니다."
        title = slide.shapes.add_textbox(Emu(914400), Emu(457200), Emu(4572000), Emu(914400))
        title.text_frame.text = "제목입니다."

    pages = parse_document("order.pptx", _pptx_bytes(build))
    assert pages[0].text.index("제목입니다.") < pages[0].text.index("본문 문장입니다.")


def test_pptx_without_text_is_rejected() -> None:
    def build(presentation) -> None:  # noqa: ANN001
        presentation.slides.add_slide(presentation.slide_layouts[6])

    with pytest.raises(DocumentError) as exc:
        parse_document("images_only.pptx", _pptx_bytes(build))
    assert "슬라이드에서 텍스트를 찾지 못했습니다" in str(exc.value)


def test_broken_pptx_gives_korean_message() -> None:
    """확장자만 .pptx 로 바꾼 구버전 .ppt 파일이 흔하다."""
    with pytest.raises(DocumentError) as exc:
        parse_document("legacy.pptx", b"not a zip archive")
    assert "PPTX를 열 수 없습니다" in str(exc.value)


def test_legacy_ppt_extension_is_guided() -> None:
    with pytest.raises(DocumentError) as exc:
        parse_document("deck.ppt", b"content")
    assert ".pptx" in str(exc.value)


def test_unsupported_extension_is_rejected() -> None:
    with pytest.raises(DocumentError) as exc:
        parse_document("archive.zip", b"content")
    assert "PDF, PPTX, TXT" in str(exc.value)


def test_empty_file_is_rejected() -> None:
    with pytest.raises(DocumentError):
        parse_document("empty.txt", b"")


def test_whitespace_only_document_is_rejected() -> None:
    with pytest.raises(DocumentError) as exc:
        parse_document("blank.txt", "   \n\n   ".encode("utf-8"))
    assert "텍스트를 찾지 못했습니다" in str(exc.value)
